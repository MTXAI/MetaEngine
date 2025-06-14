"""
RAG 基础功能实现, 包括文件加载, 分词, 检索
"""

## Loader
import csv
import importlib
import json
from functools import lru_cache
from io import TextIOWrapper
from typing import *

import chardet
import cv2
import langchain_community
import numpy as np
from PIL import Image
import tqdm
from rapidocr_onnxruntime import RapidOCR
from langchain.docstore.document import Document
from langchain_community.document_loaders import CSVLoader as BaseCSVLoader, JSONLoader as BaseJSONLoader
from langchain_community.document_loaders.helpers import detect_file_encodings
from langchain_community.document_loaders.unstructured import UnstructuredFileLoader


LOADER_DICT = {
    "UnstructuredHTMLLoader": [".html", ".htm"],
    "MHTMLLoader": [".mhtml"],
    "TextLoader": [".md"],
    "UnstructuredMarkdownLoader": [".md"],
    "JSONLoader": [".json"],
    "JSONLinesLoader": [".jsonl"],
    "CSVLoader": [".csv"],
    # "FilteredCSVLoader": [".csv"], 如果使用自定义分割csv
    "RapidOCRPDFLoader": [".pdf"],
    "RapidOCRDocLoader": [".docx"],
    "RapidOCRPPTLoader": [
        ".ppt",
        ".pptx",
    ],
    "RapidOCRLoader": [".png", ".jpg", ".jpeg", ".bmp"],
    "UnstructuredFileLoader": [
        ".eml",
        ".msg",
        ".rst",
        ".rtf",
        ".txt",
        ".xml",
        ".epub",
        ".odt",
        ".tsv",
    ],
    "UnstructuredEmailLoader": [".eml", ".msg"],
    "UnstructuredEPubLoader": [".epub"],
    "UnstructuredExcelLoader": [".xlsx", ".xls", ".xlsd"],
    "NotebookLoader": [".ipynb"],
    "UnstructuredODTLoader": [".odt"],
    "PythonLoader": [".py"],
    "UnstructuredRSTLoader": [".rst"],
    "UnstructuredRTFLoader": [".rtf"],
    "SRTLoader": [".srt"],
    "TomlLoader": [".toml"],
    "UnstructuredTSVLoader": [".tsv"],
    "UnstructuredWordDocumentLoader": [".docx"],
    "UnstructuredXMLLoader": [".xml"],
    "UnstructuredPowerPointLoader": [".ppt", ".pptx"],
    "EverNoteLoader": [".enex"],
}
SUPPORTED_EXTS = [ext for sublist in LOADER_DICT.values() for ext in sublist]


# patch json.dumps to disable ensure_ascii
def _new_json_dumps(obj, **kwargs):
    kwargs["ensure_ascii"] = False
    return _origin_json_dumps(obj, **kwargs)


if json.dumps is not _new_json_dumps:
    _origin_json_dumps = json.dumps
    json.dumps = _new_json_dumps


class JSONLoader(BaseJSONLoader):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._json_lines = True


langchain_community.document_loaders.JSONLoader = JSONLoader


class CSVLoader(BaseCSVLoader):
    def __init__(
        self,
        file_path: str,
        columns_to_read: List[str] = [],
        source_column: Optional[str] = None,
        metadata_columns: List[str] = [],
        csv_args: Optional[Dict] = None,
        encoding: Optional[str] = None,
        autodetect_encoding: bool = False,
    ):
        super().__init__(
            file_path=file_path,
            source_column=source_column,
            metadata_columns=metadata_columns,
            csv_args=csv_args,
            encoding=encoding,
            autodetect_encoding=autodetect_encoding,
        )
        self.columns_to_read = columns_to_read

    def load(self) -> List[Document]:
        """Load data into document objects."""

        docs = []
        try:
            with open(self.file_path, newline="", encoding=self.encoding) as csvfile:
                docs = self.__read_file(csvfile)
        except UnicodeDecodeError as e:
            if self.autodetect_encoding:
                detected_encodings = detect_file_encodings(self.file_path)
                for encoding in detected_encodings:
                    try:
                        with open(
                            self.file_path, newline="", encoding=encoding.encoding
                        ) as csvfile:
                            docs = self.__read_file(csvfile)
                            break
                    except UnicodeDecodeError:
                        continue
            else:
                raise RuntimeError(f"Error loading {self.file_path}") from e
        except Exception as e:
            raise RuntimeError(f"Error loading {self.file_path}") from e

        return docs

    def __read_file(self, csvfile: TextIOWrapper) -> List[Document]:
        docs = []
        csv_reader = csv.DictReader(csvfile, **self.csv_args)  # type: ignore
        if len(self.columns_to_read) == 0:
            self.columns_to_read = csv_reader.fieldnames
        for i, row in enumerate(csv_reader):
            content = []
            for col in self.columns_to_read:
                if col in row:
                    content.append(f"{col}:{str(row[col])}")
                else:
                    raise ValueError(
                        f"Column '{self.columns_to_read[0]}' not found in CSV file."
                    )
            content = "\n".join(content)
            # Extract the source if available
            source = (
                row.get(self.source_column, None)
                if self.source_column is not None
                else self.file_path
            )
            metadata = {"source": source, "row": i}

            for col in self.metadata_columns:
                if col in row:
                    metadata[col] = row[col]

            doc = Document(page_content=content, metadata=metadata)
            docs.append(doc)

        return docs


class OCRDocLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        def doc2text(filepath):
            from io import BytesIO

            import numpy as np
            from docx import Document, ImagePart
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import Table, _Cell
            from docx.text.paragraph import Paragraph
            from PIL import Image
            from rapidocr_onnxruntime import RapidOCR

            ocr = RapidOCR()
            doc = Document(filepath)
            resp = ""

            def iter_block_items(parent):
                from docx.document import Document

                if isinstance(parent, Document):
                    parent_elm = parent.element.body
                elif isinstance(parent, _Cell):
                    parent_elm = parent._tc
                else:
                    raise ValueError("RapidOCRDocLoader parse fail")

                for child in parent_elm.iterchildren():
                    if isinstance(child, CT_P):
                        yield Paragraph(child, parent)
                    elif isinstance(child, CT_Tbl):
                        yield Table(child, parent)

            b_unit = tqdm.tqdm(
                total=len(doc.paragraphs) + len(doc.tables),
                desc="RapidOCRDocLoader block index: 0",
            )
            for i, block in enumerate(iter_block_items(doc)):
                b_unit.set_description("RapidOCRDocLoader  block index: {}".format(i))
                b_unit.refresh()
                if isinstance(block, Paragraph):
                    resp += block.text.strip() + "\n"
                    images = block._element.xpath(".//pic:pic")  # 获取所有图片
                    for image in images:
                        for img_id in image.xpath(".//a:blip/@r:embed"):  # 获取图片id
                            part = doc.part.related_parts[
                                img_id
                            ]  # 根据图片id获取对应的图片
                            if isinstance(part, ImagePart):
                                image = Image.open(BytesIO(part._blob))
                                result, _ = ocr(np.array(image))
                                if result:
                                    ocr_result = [line[1] for line in result]
                                    resp += "\n".join(ocr_result)
                elif isinstance(block, Table):
                    for row in block.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                b_unit.update(1)
            return resp

        text = doc2text(self.file_path)
        from unstructured.partition.text import partition_text

        return partition_text(text=text, **self.unstructured_kwargs)


class OCRLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        def img2text(filepath):
            resp = ""
            ocr = RapidOCR()
            result, _ = ocr(filepath)
            if result:
                ocr_result = [line[1] for line in result]
                resp += "\n".join(ocr_result)
            return resp

        text = img2text(self.file_path)
        from unstructured.partition.text import partition_text

        return partition_text(text=text, **self.unstructured_kwargs)


class OCRPDFLoader(UnstructuredFileLoader):
    def __init__(self, file_path, ocr_threshold=(0.6, 0.6)):
        super().__init__(file_path)
        self.ocr_threshold = ocr_threshold

    def _get_elements(self) -> List:
        def rotate_img(img, angle):
            """
            img   --image
            angle --rotation angle
            return--rotated img
            """

            h, w = img.shape[:2]
            rotate_center = (w / 2, h / 2)
            # 获取旋转矩阵
            # 参数1为旋转中心点;
            # 参数2为旋转角度,正值-逆时针旋转;负值-顺时针旋转
            # 参数3为各向同性的比例因子,1.0原图，2.0变成原来的2倍，0.5变成原来的0.5倍
            M = cv2.getRotationMatrix2D(rotate_center, angle, 1.0)
            # 计算图像新边界
            new_w = int(h * np.abs(M[0, 1]) + w * np.abs(M[0, 0]))
            new_h = int(h * np.abs(M[0, 0]) + w * np.abs(M[0, 1]))
            # 调整旋转矩阵以考虑平移
            M[0, 2] += (new_w - w) / 2
            M[1, 2] += (new_h - h) / 2

            rotated_img = cv2.warpAffine(img, M, (new_w, new_h))
            return rotated_img

        def pdf2text(filepath):
            import fitz  # pyMuPDF里面的fitz包，不要与pip install fitz混淆
            import numpy as np

            ocr = RapidOCR()
            doc = fitz.open(filepath)
            resp = ""

            b_unit = tqdm.tqdm(
                total=doc.page_count, desc="RapidOCRPDFLoader context page index: 0"
            )
            for i, page in enumerate(doc):
                b_unit.set_description(
                    "RapidOCRPDFLoader context page index: {}".format(i)
                )
                b_unit.refresh()
                text = page.get_text("")
                resp += text + "\n"

                img_list = page.get_image_info(xrefs=True)
                for img in img_list:
                    if xref := img.get("xref"):
                        bbox = img["bbox"]
                        # 检查图片尺寸是否超过设定的阈值
                        if (bbox[2] - bbox[0]) / (page.rect.width) < self.ocr_threshold[
                            0
                        ] or (bbox[3] - bbox[1]) / (
                            page.rect.height
                        ) < self.ocr_threshold[1]:
                            continue
                        pix = fitz.Pixmap(doc, xref)
                        samples = pix.samples
                        if int(page.rotation) != 0:  # 如果Page有旋转角度，则旋转图片
                            img_array = np.frombuffer(
                                pix.samples, dtype=np.uint8
                            ).reshape(pix.height, pix.width, -1)
                            tmp_img = Image.fromarray(img_array)
                            ori_img = cv2.cvtColor(np.array(tmp_img), cv2.COLOR_RGB2BGR)
                            rot_img = rotate_img(img=ori_img, angle=360 - page.rotation)
                            img_array = cv2.cvtColor(rot_img, cv2.COLOR_RGB2BGR)
                        else:
                            img_array = np.frombuffer(
                                pix.samples, dtype=np.uint8
                            ).reshape(pix.height, pix.width, -1)

                        result, _ = ocr(img_array)
                        if result:
                            ocr_result = [line[1] for line in result]
                            resp += "\n".join(ocr_result)

                # 更新进度
                b_unit.update(1)
            return resp

        text = pdf2text(self.file_path)
        from unstructured.partition.text import partition_text

        return partition_text(text=text, **self.unstructured_kwargs)


class OCRPPTLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        def ppt2text(filepath):
            from io import BytesIO

            import numpy as np
            from PIL import Image
            from pptx import Presentation
            from rapidocr_onnxruntime import RapidOCR

            ocr = RapidOCR()
            prs = Presentation(filepath)
            resp = ""

            def extract_text(shape):
                nonlocal resp
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                if shape.shape_type == 13:  # 13 表示图片
                    image = Image.open(BytesIO(shape.image.blob))
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                elif shape.shape_type == 6:  # 6 表示组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)

            b_unit = tqdm.tqdm(
                total=len(prs.slides), desc="RapidOCRPPTLoader slide index: 1"
            )
            # 遍历所有幻灯片
            for slide_number, slide in enumerate(prs.slides, start=1):
                b_unit.set_description(
                    "RapidOCRPPTLoader slide index: {}".format(slide_number)
                )
                b_unit.refresh()
                sorted_shapes = sorted(
                    slide.shapes, key=lambda x: (x.top, x.left)
                )  # 从上到下、从左到右遍历
                for shape in sorted_shapes:
                    extract_text(shape)
                b_unit.update(1)
            return resp

        text = ppt2text(self.file_path)
        from unstructured.partition.text import partition_text

        return partition_text(text=text, **self.unstructured_kwargs)


def get_LoaderClass(file_extension):
    for LoaderClass, extensions in LOADER_DICT.items():
        if file_extension in extensions:
            return LoaderClass


def get_loader(loader_name: str, file_path: str, loader_kwargs: Dict = None):
    """
    根据loader_name和文件路径或内容返回文档加载器。
    """
    loader_kwargs = loader_kwargs or {}
    try:
        if loader_name in [
            "RapidOCRPDFLoader",
            "RapidOCRLoader",
            "FilteredCSVLoader",
            "RapidOCRDocLoader",
            "RapidOCRPPTLoader",
        ]:
            document_loaders_module = importlib.import_module(
                "chatchat.server.file_rag.document_loaders"
            )
        else:
            document_loaders_module = importlib.import_module(
                "langchain_community.document_loaders"
            )
        DocumentLoader = getattr(document_loaders_module, loader_name)
    except Exception as e:
        msg = f"为文件{file_path}查找加载器{loader_name}时出错：{e}"
        logger.error(f"{e.__class__.__name__}: {msg}")
        document_loaders_module = importlib.import_module(
            "langchain_community.document_loaders"
        )
        DocumentLoader = getattr(document_loaders_module, "UnstructuredFileLoader")

    if loader_name == "UnstructuredFileLoader":
        loader_kwargs.setdefault("autodetect_encoding", True)
    elif loader_name == "CSVLoader":
        if not loader_kwargs.get("encoding"):
            # 如果未指定 encoding，自动识别文件编码类型，避免langchain loader 加载文件报编码错误
            with open(file_path, "rb") as struct_file:
                encode_detect = chardet.detect(struct_file.read())
            if encode_detect is None:
                encode_detect = {"encoding": "utf-8"}
            loader_kwargs["encoding"] = encode_detect["encoding"]

    elif loader_name == "JSONLoader":
        loader_kwargs.setdefault("jq_schema", ".")
        loader_kwargs.setdefault("text_content", False)

    loader = DocumentLoader(file_path, **loader_kwargs)
    return loader


## Spliter
import re

from langchain.text_splitter import CharacterTextSplitter, TextSplitter, MarkdownHeaderTextSplitter
from langchain.text_splitter import RecursiveCharacterTextSplitter


text_splitter_dict: Dict[str, Dict[str, Any]] = {
            "ChineseRecursiveTextSplitter": {
                "source": "",
                "tokenizer_name_or_path": "",
            },
            "SpacyTextSplitter": {
                "source": "huggingface",
                "tokenizer_name_or_path": "gpt2",
            },
            "RecursiveCharacterTextSplitter": {
                "source": "tiktoken",
                "tokenizer_name_or_path": "cl100k_base",
            },
            "MarkdownHeaderTextSplitter": {
                "headers_to_split_on": [
                    ("#", "head1"),
                    ("##", "head2"),
                    ("###", "head3"),
                    ("####", "head4"),
                ]
            },
        }


class AliTextSplitter(CharacterTextSplitter):
    def __init__(self, pdf: bool = False, **kwargs):
        super().__init__(**kwargs)
        self.pdf = pdf

    def split_text(self, text: str) -> List[str]:
        # use_document_segmentation参数指定是否用语义切分文档，此处采取的文档语义分割模型为达摩院开源的nlp_bert_document-segmentation_chinese-base，论文见https://arxiv.org/abs/2107.09278
        # 如果使用模型进行文档语义切分，那么需要安装modelscope[nlp]：pip install "modelscope[nlp]" -f https://modelscope.oss-cn-beijing.aliyuncs.com/releases/repo.html
        # 考虑到使用了三个模型，可能对于低配置gpu不太友好，因此这里将模型load进cpu计算，有需要的话可以替换device为自己的显卡id
        if self.pdf:
            text = re.sub(r"\n{3,}", r"\n", text)
            text = re.sub("\s", " ", text)
            text = re.sub("\n\n", "", text)
        try:
            from modelscope.pipelines import pipeline
        except ImportError:
            raise ImportError(
                "Could not import modelscope python package. "
                "Please install modelscope with `pip install modelscope`. "
            )

        p = pipeline(
            task="document-segmentation",
            model="damo/nlp_bert_document-segmentation_chinese-base",
            device="cpu",
        )
        result = p(documents=text)
        sent_list = [i for i in result["text"].split("\n\t") if i]
        return sent_list


def _split_text_with_regex_from_end(
    text: str, separator: str, keep_separator: bool
) -> List[str]:
    # Now that we have the separator, split the text
    if separator:
        if keep_separator:
            # The parentheses in the pattern keep the delimiters in the result.
            _splits = re.split(f"({separator})", text)
            splits = ["".join(i) for i in zip(_splits[0::2], _splits[1::2])]
            if len(_splits) % 2 == 1:
                splits += _splits[-1:]
            # splits = [_splits[0]] + splits
        else:
            splits = re.split(separator, text)
    else:
        splits = list(text)
    return [s for s in splits if s != ""]


class ChineseRecursiveTextSplitter(RecursiveCharacterTextSplitter):
    def __init__(
        self,
        separators: Optional[List[str]] = None,
        keep_separator: bool = True,
        is_separator_regex: bool = True,
        **kwargs: Any,
    ) -> None:
        """Create a new TextSplitter."""
        super().__init__(keep_separator=keep_separator, **kwargs)
        self._separators = separators or [
            "\n\n",
            "\n",
            "。|！|？",
            "\.\s|\!\s|\?\s",
            "；|;\s",
            "，|,\s",
        ]
        self._is_separator_regex = is_separator_regex

    def _split_text(self, text: str, separators: List[str]) -> List[str]:
        """Split incoming text and return chunks."""
        final_chunks = []
        # Get appropriate separator to use
        separator = separators[-1]
        new_separators = []
        for i, _s in enumerate(separators):
            _separator = _s if self._is_separator_regex else re.escape(_s)
            if _s == "":
                separator = _s
                break
            if re.search(_separator, text):
                separator = _s
                new_separators = separators[i + 1 :]
                break

        _separator = separator if self._is_separator_regex else re.escape(separator)
        splits = _split_text_with_regex_from_end(text, _separator, self._keep_separator)

        # Now go merging things, recursively splitting longer texts.
        _good_splits = []
        _separator = "" if self._keep_separator else separator
        for s in splits:
            if self._length_function(s) < self._chunk_size:
                _good_splits.append(s)
            else:
                if _good_splits:
                    merged_text = self._merge_splits(_good_splits, _separator)
                    final_chunks.extend(merged_text)
                    _good_splits = []
                if not new_separators:
                    final_chunks.append(s)
                else:
                    other_info = self._split_text(s, new_separators)
                    final_chunks.extend(other_info)
        if _good_splits:
            merged_text = self._merge_splits(_good_splits, _separator)
            final_chunks.extend(merged_text)
        return [
            re.sub(r"\n{2,}", "\n", chunk.strip())
            for chunk in final_chunks
            if chunk.strip() != ""
        ]


class ChineseTextSplitter(CharacterTextSplitter):
    def __init__(self, pdf: bool = False, sentence_size: int = 250, **kwargs):
        super().__init__(**kwargs)
        self.pdf = pdf
        self.sentence_size = sentence_size

    def split_text1(self, text: str) -> List[str]:
        if self.pdf:
            text = re.sub(r"\n{3,}", "\n", text)
            text = re.sub("\s", " ", text)
            text = text.replace("\n\n", "")
        sent_sep_pattern = re.compile(
            '([﹒﹔﹖﹗．。！？]["’”」』]{0,2}|(?=["‘“「『]{1,2}|$))'
        )  # del ：；
        sent_list = []
        for ele in sent_sep_pattern.split(text):
            if sent_sep_pattern.match(ele) and sent_list:
                sent_list[-1] += ele
            elif ele:
                sent_list.append(ele)
        return sent_list

    def split_text(self, text: str) -> List[str]:  ##此处需要进一步优化逻辑
        if self.pdf:
            text = re.sub(r"\n{3,}", r"\n", text)
            text = re.sub("\s", " ", text)
            text = re.sub("\n\n", "", text)

        text = re.sub(r"([;；.!?。！？\?])([^”’])", r"\1\n\2", text)  # 单字符断句符
        text = re.sub(r'(\.{6})([^"’”」』])', r"\1\n\2", text)  # 英文省略号
        text = re.sub(r'(\…{2})([^"’”」』])', r"\1\n\2", text)  # 中文省略号
        text = re.sub(
            r'([;；!?。！？\?]["’”」』]{0,2})([^;；!?，。！？\?])', r"\1\n\2", text
        )
        # 如果双引号前有终止符，那么双引号才是句子的终点，把分句符\n放到双引号后，注意前面的几句都小心保留了双引号
        text = text.rstrip()  # 段尾如果有多余的\n就去掉它
        # 很多规则中会考虑分号;，但是这里我把它忽略不计，破折号、英文双引号等同样忽略，需要的再做些简单调整即可。
        ls = [i for i in text.split("\n") if i]
        for ele in ls:
            if len(ele) > self.sentence_size:
                ele1 = re.sub(r'([,，.]["’”」』]{0,2})([^,，.])', r"\1\n\2", ele)
                ele1_ls = ele1.split("\n")
                for ele_ele1 in ele1_ls:
                    if len(ele_ele1) > self.sentence_size:
                        ele_ele2 = re.sub(
                            r'([\n]{1,}| {2,}["’”」』]{0,2})([^\s])',
                            r"\1\n\2",
                            ele_ele1,
                        )
                        ele2_ls = ele_ele2.split("\n")
                        for ele_ele2 in ele2_ls:
                            if len(ele_ele2) > self.sentence_size:
                                ele_ele3 = re.sub(
                                    '( ["’”」』]{0,2})([^ ])', r"\1\n\2", ele_ele2
                                )
                                ele2_id = ele2_ls.index(ele_ele2)
                                ele2_ls = (
                                    ele2_ls[:ele2_id]
                                    + [i for i in ele_ele3.split("\n") if i]
                                    + ele2_ls[ele2_id + 1 :]
                                )
                        ele_id = ele1_ls.index(ele_ele1)
                        ele1_ls = (
                            ele1_ls[:ele_id]
                            + [i for i in ele2_ls if i]
                            + ele1_ls[ele_id + 1 :]
                        )

                id = ls.index(ele)
                ls = ls[:id] + [i for i in ele1_ls if i] + ls[id + 1 :]
        return ls


def _under_non_alpha_ratio(text: str, threshold: float = 0.5):
    """Checks if the proportion of non-alpha characters in the text snippet exceeds a given
    threshold. This helps prevent text like "-----------BREAK---------" from being tagged
    as a title or narrative text. The ratio does not count spaces.

    Parameters
    ----------
    text
        The input string to test
    threshold
        If the proportion of non-alpha characters exceeds this threshold, the function
        returns False
    """
    if len(text) == 0:
        return False

    alpha_count = len([char for char in text if char.strip() and char.isalpha()])
    total_count = len([char for char in text if char.strip()])
    try:
        ratio = alpha_count / total_count
        return ratio < threshold
    except:
        return False


def _is_possible_title(
    text: str,
    title_max_word_length: int = 20,
    non_alpha_threshold: float = 0.5,
) -> bool:
    """Checks to see if the text passes all of the checks for a valid title.

    Parameters
    ----------
    text
        The input text to check
    title_max_word_length
        The maximum number of words a title can contain
    non_alpha_threshold
        The minimum number of alpha characters the text needs to be considered a title
    """

    # 文本长度为0的话，肯定不是title
    if len(text) == 0:
        print("Not a title. Text is empty.")
        return False

    # 文本中有标点符号，就不是title
    ENDS_IN_PUNCT_PATTERN = r"[^\w\s]\Z"
    ENDS_IN_PUNCT_RE = re.compile(ENDS_IN_PUNCT_PATTERN)
    if ENDS_IN_PUNCT_RE.search(text) is not None:
        return False

    # 文本长度不能超过设定值，默认20
    # NOTE(robinson) - splitting on spaces here instead of word tokenizing because it
    # is less expensive and actual tokenization doesn't add much value for the length check
    if len(text) > title_max_word_length:
        return False

    # 文本中数字的占比不能太高，否则不是title
    if _under_non_alpha_ratio(text, threshold=non_alpha_threshold):
        return False

    # NOTE(robinson) - Prevent flagging salutations like "To My Dearest Friends," as titles
    if text.endswith((",", ".", "，", "。")):
        return False

    if text.isnumeric():
        print(f"Not a title. Text is all numeric:\n\n{text}")  # type: ignore
        return False

    # 开头的字符内应该有数字，默认5个字符内
    if len(text) < 5:
        text_5 = text
    else:
        text_5 = text[:5]
    alpha_in_text_5 = sum(list(map(lambda x: x.isnumeric(), list(text_5))))
    if not alpha_in_text_5:
        return False

    return True


def zh_title_enhance(docs: Document) -> Optional[Document]:
    title = None
    if len(docs) > 0:
        for doc in docs:
            if _is_possible_title(doc.page_content):
                doc.metadata["category"] = "cn_Title"
                title = doc.page_content
            elif title:
                doc.page_content = f"下文与({title})有关。{doc.page_content}"
        return docs
    else:
        return None


@lru_cache()
def make_text_splitter(splitter_name, chunk_size, chunk_overlap):
    """
    根据参数获取特定的分词器
    """
    splitter_name = splitter_name or "SpacyTextSplitter"
    try:
        if (
            splitter_name == "MarkdownHeaderTextSplitter"
        ):  # MarkdownHeaderTextSplitter特殊判定
            headers_to_split_on = text_splitter_dict[splitter_name][
                "headers_to_split_on"
            ]
            text_splitter = MarkdownHeaderTextSplitter(
                headers_to_split_on=headers_to_split_on, strip_headers=False
            )
        else:
            try:  # 优先使用用户自定义的text_splitter
                text_splitter_module = importlib.import_module("chatchat.server.file_rag.text_splitter")
                TextSplitter = getattr(text_splitter_module, splitter_name)
            except:  # 否则使用langchain的text_splitter
                text_splitter_module = importlib.import_module(
                    "langchain.text_splitter"
                )
                TextSplitter = getattr(text_splitter_module, splitter_name)

            if (
                text_splitter_dict[splitter_name]["source"] == "tiktoken"
            ):  # 从tiktoken加载
                try:
                    text_splitter = TextSplitter.from_tiktoken_encoder(
                        encoding_name=text_splitter_dict[splitter_name][
                            "tokenizer_name_or_path"
                        ],
                        pipeline="zh_core_web_sm",
                        chunk_size=chunk_size,
                        chunk_overlap=chunk_overlap,
                    )
                except:
                    text_splitter = TextSplitter.from_tiktoken_encoder(
                        encoding_name=text_splitter_dict[splitter_name][
                            "tokenizer_name_or_path"
                        ],
                        chunk_size=chunk_size,
                        chunk_overlap=chunk_overlap,
                    )
            elif (
                text_splitter_dict[splitter_name]["source"] == "huggingface"
            ):  # 从huggingface加载
                if (
                    text_splitter_dict[splitter_name]["tokenizer_name_or_path"]
                    == "gpt2"
                ):
                    from langchain.text_splitter import CharacterTextSplitter
                    from transformers import GPT2TokenizerFast

                    tokenizer = GPT2TokenizerFast.from_pretrained("gpt2")
                else:  # 字符长度加载
                    from transformers import AutoTokenizer

                    tokenizer = AutoTokenizer.from_pretrained(
                        text_splitter_dict[splitter_name]["tokenizer_name_or_path"],
                        trust_remote_code=True,
                    )
                text_splitter = TextSplitter.from_huggingface_tokenizer(
                    tokenizer=tokenizer,
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                )
            else:
                try:
                    text_splitter = TextSplitter(
                        pipeline="zh_core_web_sm",
                        chunk_size=chunk_size,
                        chunk_overlap=chunk_overlap,
                    )
                except:
                    text_splitter = TextSplitter(
                        chunk_size=chunk_size, chunk_overlap=chunk_overlap
                    )
    except Exception as e:
        print(e)
        text_splitter_module = importlib.import_module("langchain.text_splitter")
        TextSplitter = getattr(text_splitter_module, "RecursiveCharacterTextSplitter")
        text_splitter = TextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)

    # If you use SpacyTextSplitter you can use GPU to do split likes Issue #1287
    # text_splitter._tokenizer.max_length = 37016792
    # text_splitter._tokenizer.prefer_gpu()
    return text_splitter


## Retriever
from abc import ABCMeta, abstractmethod

from langchain.retrievers import EnsembleRetriever
from langchain.vectorstores import VectorStore
from langchain_community.retrievers import BM25Retriever
from langchain_core.retrievers import BaseRetriever
from langchain.docstore.document import Document


class BaseRetrieverService(metaclass=ABCMeta):
    def __init__(self, **kwargs):
        self.do_init(**kwargs)

    @abstractmethod
    def do_init(self, **kwargs):
        pass

    @abstractmethod
    def from_vectorstore(
        vectorstore: VectorStore,
        top_k: int,
        score_threshold: int | float,
    ):
        pass

    @abstractmethod
    def get_relevant_documents(self, query: str):
        pass

# for faiss
class EnsembleRetrieverService(BaseRetrieverService):
    def do_init(
        self,
        retriever: BaseRetriever = None,
        top_k: int = 5,
    ):
        self.vs = None
        self.top_k = top_k
        self.retriever = retriever

    @staticmethod
    def from_vectorstore(
        vectorstore: VectorStore,
        top_k: int,
        score_threshold: int | float,
    ):
        faiss_retriever = vectorstore.as_retriever(
            search_type="similarity_score_threshold",
            search_kwargs={"score_threshold": score_threshold, "k": top_k},
        )
        # TODO: 换个不用torch的实现方式
        # from cutword.cutword import Cutter
        import jieba

        # cutter = Cutter()
        docs = list(vectorstore.docstore._dict.values())
        bm25_retriever = BM25Retriever.from_documents(
            docs,
            preprocess_func=jieba.lcut_for_search,
        )
        bm25_retriever.k = top_k
        ensemble_retriever = EnsembleRetriever(
            retrievers=[bm25_retriever, faiss_retriever], weights=[0.5, 0.5]
        )
        return EnsembleRetrieverService(retriever=ensemble_retriever, top_k=top_k)

    def get_relevant_documents(self, query: str):
        return self.retriever.get_relevant_documents(query)[: self.top_k]

# for vector store
class VectorstoreRetrieverService(BaseRetrieverService):
    def do_init(
        self,
        retriever: BaseRetriever = None,
        top_k: int = 5,
    ):
        self.vs = None
        self.top_k = top_k
        self.retriever = retriever

    @staticmethod
    def from_vectorstore(
        vectorstore: VectorStore,
        top_k: int,
        score_threshold: int | float,
    ):
        retriever = vectorstore.as_retriever(
            search_type="similarity_score_threshold",
            search_kwargs={"score_threshold": score_threshold, "k": top_k},
        )
        return VectorstoreRetrieverService(retriever=retriever, top_k=top_k)

    def get_relevant_documents(self, query: str):
        return self.retriever.get_relevant_documents(query)[: self.top_k]

